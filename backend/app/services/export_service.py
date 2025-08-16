# app/services/export_service.py
from __future__ import annotations

from pathlib import Path
from datetime import datetime
import io
import re
import tempfile
import urllib.request
from typing import Iterable, Optional

from app.core.config import settings
from app.core.telemetry import aspan
from app.models.schemas.slide import Slide
from app.models.schemas.export import (
    ExportResponse,
    EditorDocIn,
    EditorSlideIn,
    EditorLayer,
)

_SLIDE_PREFIX = re.compile(r"^\s*Slide\s+\d+\s*:\s*", re.IGNORECASE)

# PowerPoint uses English Metric Units (EMU); 1 inch = 914,400 EMU; CSS assumes 96 px = 1 inch
EMU_PER_INCH = 914_400
DPI = 96
EMU_PER_PX = int(EMU_PER_INCH / DPI)  # 9,525


def _emu(px: float) -> int:
    return int(round(px * EMU_PER_PX))


def _pt_from_px(px: float) -> float:
    # 1 pt = 1/72 inch; 96 px = 72 pt  =>  px * 0.75
    return float(px) * 0.75


def _strip_slide_prefix(s: str) -> str:
    return _SLIDE_PREFIX.sub("", (s or "").strip())


def _export_dir() -> Path:
    d = (Path(settings.STORAGE_DIR) / "exports").resolve()
    d.mkdir(parents=True, exist_ok=True)
    return d


def _stamp_name(theme: str) -> str:
    return f"deck_{datetime.utcnow().strftime('%Y%m%d_%H%M%S')}_{theme}"


def _fetch_image_png_bytes(url: str) -> Optional[bytes]:
    """
    Fetch image bytes from URL and normalize to PNG using Pillow.
    Returns PNG bytes (or original bytes if Pillow missing) or None on failure.
    """
    if not url:
        return None
    try:
        req = urllib.request.Request(
            url,
            headers={
                "User-Agent": "PresenTuneAI/1.0 (+https://example)",
                "Accept": "image/*,*/*;q=0.8",
            },
        )
        with urllib.request.urlopen(req, timeout=15) as resp:
            data = resp.read()
    except Exception:
        return None

    # Try Pillow normalization (handles WEBP, etc.)
    try:
        from PIL import Image  # pillow is a transitive dep of python-pptx
        with Image.open(io.BytesIO(data)) as im:
            if im.mode not in ("RGB", "RGBA"):
                im = im.convert("RGBA" if "A" in im.getbands() else "RGB")
            out = io.BytesIO()
            im.save(out, format="PNG", optimize=True)
            return out.getvalue()
    except Exception:
        # If Pillow not available or decode fails, return original bytes;
        # python-pptx supports PNG/JPEG/GIF/BMP/TIFF and may still succeed.
        return data or None


# ---------- SLIDES (old/simple) ----------
async def _export_slides_to_pptx(slides: list[Slide], theme: str, out_path: Path) -> ExportResponse:
    try:
        from pptx import Presentation
        from pptx.util import Pt, Inches
    except Exception as e:
        # txt fallback
        txt = out_path.with_suffix(".txt")
        async with aspan(
            "export_txt_fallback",
            theme=theme,
            slides=len(slides),
            out=str(txt),
            reason=type(e).__name__,
        ):
            with txt.open("w", encoding="utf-8") as f:
                for idx, s in enumerate(slides, start=1):
                    title = _strip_slide_prefix(s.title or f"Slide {idx}")
                    f.write(f"{title}\n")
                    for b in (s.bullets or []):
                        f.write(f"  - {b}\n")
                    if getattr(s, "notes", None):
                        f.write(f"  [notes] {s.notes}\n")
                    for m in (s.media or []):
                        if isinstance(m, dict):
                            url = m.get("url"); alt = m.get("alt"); typ = m.get("type", "image")
                        else:
                            url = getattr(m, "url", None); alt = getattr(m, "alt", None); typ = getattr(m, "type", "image")
                        f.write(f"  [media] {typ} {url or ''}" + (f" — {alt}" if alt else "") + "\n")
                    f.write("\n")
        return ExportResponse(path=str(txt), format="txt", theme=theme, bytes=txt.stat().st_size)

    async with aspan("export_pptx_simple", theme=theme, slides=len(slides), out=str(out_path)):
        prs = Presentation()
        # Prefer Title + Content (1). If unavailable, fall back to Title (0).
        content_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]

        for idx, s in enumerate(slides, start=1):
            slide = prs.slides.add_slide(content_layout)
            title_text = _strip_slide_prefix(s.title or f"Slide {idx}")

            # Title
            if slide.shapes.title:
                slide.shapes.title.text = title_text
                tf = slide.shapes.title.text_frame
                if tf.paragraphs:
                    p = tf.paragraphs[0]
                    if p.runs:
                        for r in p.runs:
                            r.font.size = Pt(32)
                    else:
                        p.font.size = Pt(32)

            # Bullets
            body = slide.placeholders[1].text_frame if len(slide.placeholders) > 1 else None
            if body:
                body.clear()
                first = True
                for b in (s.bullets or []):
                    if first:
                        body.text = b
                        if body.paragraphs:
                            body.paragraphs[0].font.size = Pt(18)
                        first = False
                    else:
                        p = body.add_paragraph()
                        p.text = b
                        p.level = 0
                        p.font.size = Pt(18)

            # First image (best effort) on the right
            img_url = None
            for m in (s.media or []):
                url = (m.get("url") if isinstance(m, dict) else getattr(m, "url", None))
                if url:
                    img_url = url
                    break
            if img_url:
                img_bytes = _fetch_image_png_bytes(img_url)
                if img_bytes:
                    try:
                        left = Inches(8.2)
                        top = Inches(2.0)
                        width = Inches(3.5)
                        slide.shapes.add_picture(io.BytesIO(img_bytes), left, top, width=width)
                    except Exception:
                        pass  # keep exporting other slides

            # Notes
            if getattr(s, "notes", None):
                slide.notes_slide.notes_text_frame.text = s.notes

        prs.save(str(out_path))

    return ExportResponse(path=str(out_path), format="pptx", theme=theme, bytes=out_path.stat().st_size)


# ---------- EDITOR (framed layers → exact positioning) ----------
def _find_blank_layout(prs) -> int:
    # Try to find layout named 'Blank'; else last layout
    for i, layout in enumerate(prs.slide_layouts):
        if getattr(layout, "name", "").lower() == "blank":
            return i
    return len(prs.slide_layouts) - 1


def _pp_align_from(style_align: str | None):
    # Map CSS-ish to pptx.enum.text.PP_ALIGN
    try:
        from pptx.enum.text import PP_ALIGN
        a = (style_align or "").lower()
        if a in ("center", "middle"):
            return PP_ALIGN.CENTER
        if a in ("right", "end"):
            return PP_ALIGN.RIGHT
        if a in ("justify",):
            return PP_ALIGN.JUSTIFY
        return PP_ALIGN.LEFT
    except Exception:
        return None


def _add_text_layer(slide, ly: EditorLayer, default_font: str = "Inter"):
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    f = ly.frame
    left, top, width, height = _emu(f.x), _emu(f.y), _emu(f.w), _emu(f.h)

    shp = slide.shapes.add_textbox(left, top, width, height)
    tf = shp.text_frame
    tf.clear()
    tf.word_wrap = True

    text = (ly.text or "")
    # preserve "pre-wrap": one paragraph per newline
    lines = text.splitlines() or [""]
    style = ly.style or {}
    font_name = style.get("font") or default_font
    weight = int(style.get("weight") or 400)
    size_px = float(style.get("size") or 20)
    color_hex = (style.get("color") or "#111111").lstrip("#")
    try:
        rgb = RGBColor(int(color_hex[0:2], 16), int(color_hex[2:4], 16), int(color_hex[4:6], 16))
    except Exception:
        rgb = RGBColor(17, 17, 17)

    align = _pp_align_from(style.get("align") or style.get("textAlign"))

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        if align is not None:
            p.alignment = align

        # Apply font to the whole paragraph (single-run unless rich text present)
        if p.runs:
            for r in p.runs:
                r.font.size = Pt(_pt_from_px(size_px))
                r.font.bold = bool(weight >= 600)
                r.font.name = font_name
                r.font.color.rgb = rgb
        else:
            p.font.size = Pt(_pt_from_px(size_px))
            p.font.bold = bool(weight >= 600)
            p.font.name = font_name
            p.font.color.rgb = rgb


def _add_image_layer(slide, ly: EditorLayer):
    """
    Places the image inside ly.frame honoring object-fit.
    - cover: center-crop to fill the frame (no distortion)
    - contain: scale down to fit and center (letterbox)
    - fill: stretch to frame (distort)
    """
    f = ly.frame
    left, top, width, height = _emu(f.x), _emu(f.y), _emu(f.w), _emu(f.h)
    fit = (ly.fit or "cover").lower()
    src = ly.source or {}
    url = src.get("url") if isinstance(src, dict) else getattr(src, "url", None)
    if not url:
        return

    img_bytes = _fetch_image_png_bytes(url)
    if not img_bytes:
        return

    if fit == "fill":
        slide.shapes.add_picture(io.BytesIO(img_bytes), left, top, width=width, height=height)
        return

    # Need original image size (use Pillow on normalized bytes)
    try:
        from PIL import Image
        with Image.open(io.BytesIO(img_bytes)) as im:
            iw, ih = im.size
    except Exception:
        # fallback: just stretch
        slide.shapes.add_picture(io.BytesIO(img_bytes), left, top, width=width, height=height)
        return

    # Target in px (doc frames are px)
    tw, th = float(f.w), float(f.h)
    rx, ry = tw / iw, th / ih

    if fit == "contain":
        # scale to fit within frame, center it
        scale = min(rx, ry)
        dw, dh = iw * scale, ih * scale
        cx = _emu(f.x + (tw - dw) / 2.0)
        cy = _emu(f.y + (th - dh) / 2.0)
        slide.shapes.add_picture(io.BytesIO(img_bytes), cx, cy, width=_emu(dw), height=_emu(dh))
        return

    # cover (default): crop equally from overflow side(s), keep center
    pic = slide.shapes.add_picture(io.BytesIO(img_bytes), left, top, width=width, height=height)

    # Fractions relative to original image dimension
    if rx > ry:
        # overflow in height
        keep_h_frac = max(0.0, min(1.0, th / (ih * rx)))
        crop_each = max(0.0, (1.0 - keep_h_frac) / 2.0)
        pic.crop_top = crop_each
        pic.crop_bottom = crop_each
        pic.crop_left = 0.0
        pic.crop_right = 0.0
    else:
        # overflow in width
        keep_w_frac = max(0.0, min(1.0, tw / (iw * ry)))
        crop_each = max(0.0, (1.0 - keep_w_frac) / 2.0)
        pic.crop_left = crop_each
        pic.crop_right = crop_each
        pic.crop_top = 0.0
        pic.crop_bottom = 0.0


async def _export_editor_to_pptx(doc: EditorDocIn, theme: str, out_path: Path) -> ExportResponse:
    from pptx import Presentation
    from pptx.dml.color import RGBColor

    slides = doc.slides or []
    page = doc.page or {}
    page_w = float(page.get("width", 1280))
    page_h = float(page.get("height", 720))

    async with aspan("export_pptx_editor", theme=theme, slides=len(slides), out=str(out_path)):
        prs = Presentation()

        # Match editor canvas size (px → EMU)
        prs.slide_width = _emu(page_w)
        prs.slide_height = _emu(page_h)

        blank_idx = _find_blank_layout(prs)
        for s in slides:
            sl = prs.slides.add_slide(prs.slide_layouts[blank_idx])

            # Background fill (solid)
            bg = s.background or {}
            fill_hex = (bg.get("fill") or "#FFFFFF").lstrip("#")
            try:
                r, g, b = int(fill_hex[0:2], 16), int(fill_hex[2:4], 16), int(fill_hex[4:6], 16)
                sl.background.fill.solid()
                sl.background.fill.fore_color.rgb = RGBColor(r, g, b)
            except Exception:
                pass

            # z-order low → high
            for ly in sorted(s.layers or [], key=lambda L: (L.z or 0)):
                if ly.kind == "textbox":
                    _add_text_layer(sl, ly)
                elif ly.kind == "image":
                    _add_image_layer(sl, ly)
                else:
                    continue

        prs.save(str(out_path))

    return ExportResponse(path=str(out_path), format="pptx", theme=theme, bytes=out_path.stat().st_size)


# ---------- Public entry point ----------
async def export_to_pptx(
    slides: Optional[list[Slide]] = None,
    editor: Optional[EditorDocIn] = None,
    theme: str = "default",
) -> ExportResponse:
    """
    If `editor` is provided, export exact EditorDoc frames/layers.
    Otherwise, fall back to the simple 'slides' exporter.
    """
    out_dir = _export_dir()
    base = _stamp_name(theme)
    pptx_path = out_dir / f"{base}.pptx"

    if editor is not None:
        return await _export_editor_to_pptx(editor, theme, pptx_path)
    if slides is not None:
        return await _export_slides_to_pptx(slides, theme, pptx_path)

    # Should not happen (schema validator guards), but be defensive:
    raise ValueError("export_to_pptx requires either slides or editor")
